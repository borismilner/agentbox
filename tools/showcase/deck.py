#!/usr/bin/env python3
"""Build the sales deck that accompanies docs/showcase.md.

    python3 tools/showcase/deck.py [-o docs/agentbox-showcase.pptx]

Why a generator and not a hand-made file: the deck says the same things as the
plan, and the two have to stay the same thing. A .pptx is a zip of XML that
nobody will diff, so the source of truth is this file - readable, reviewable, and
regenerable after the pitch changes.

It needs python-pptx (`pip install python-pptx`) and produces plain OOXML that
OnlyOffice, PowerPoint, Keynote and Google Slides all open. Every slide carries
its narration in the speaker notes, which tools/showcase/perform.py reads and
speaks - so the deck, the words and the live demo are one artifact.

Nothing on a slide is addressed to the presenter. No file names, no act numbers, no
"the one to sell hardest": a slide the audience can read is a slide that says only
what the audience came for. The presenter's copy of the running order is
docs/showcase.md, and it stays there.

Fonts are picked from what a Linux desktop actually has (Ubuntu Sans, JetBrains
Mono) so the file renders as designed instead of falling back to something wider
and breaking every line.
"""

import argparse
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- design system

BG = RGBColor(0x06, 0x08, 0x0A)  # near black, not black: a true black flares on a projector
BG2 = RGBColor(0x0D, 0x11, 0x17)
SURFACE = RGBColor(0x11, 0x16, 0x1D)
EDGE = RGBColor(0x1E, 0x26, 0x30)
INK = RGBColor(0xEC, 0xF1, 0xF6)
INK2 = RGBColor(0x9E, 0xAC, 0xBA)
INK3 = RGBColor(0x5E, 0x6C, 0x7A)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

SANS = "Ubuntu Sans"
MONO = "JetBrains Mono"

# The author's credit, on every slide as a watermark rather than as a slide of its
# own: the licence asks for a visible credit, and the pitch is not about him.
AUTHOR = "boris.milner@gmail.com"

W, H = 13.333, 7.5
MARGIN = 0.95
COLW = W - 2 * MARGIN

RAIN = "01アイウエオカキクケコサシスセソタチツテトナニヌネノabcdef10"


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.n = 0
        self.warnings = []
        self.credits = []

    # -- primitives ---------------------------------------------------------

    def _rect(self, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, line=None, lw=1.0):
        s = self.slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        s.shadow.inherit = False
        if fill is None:
            s.fill.background()
        else:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        if line is None:
            s.line.fill.background()
        else:
            s.line.color.rgb = line
            s.line.width = Pt(lw)
        s.text_frame.text = ""
        return s

    def _text(self, x, y, w, h, runs, size=20, color=INK, font=SANS, bold=False,
              align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP, space_after=6,
              oneline=False):
        """runs: a string, or a list of paragraphs; a paragraph is a string or a
        list of (text, {overrides}) tuples."""
        tb = self.slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        paras = runs if isinstance(runs, list) else [runs]
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            p.space_after = Pt(space_after)
            pieces = para if isinstance(para, list) else [(para, {})]
            for text, over in pieces:
                r = p.add_run()
                r.text = text
                f = r.font
                f.name = over.get("font", font)
                f.size = Pt(over.get("size", size))
                f.bold = over.get("bold", bold)
                f.italic = over.get("italic", False)
                f.color.rgb = over.get("color", color)
        if oneline:
            self._check_fits(w, paras, size, font)
        return tb

    # Rendering is the one thing python-pptx cannot do, so a box too narrow for its
    # own text stays invisible until somebody watches the slide - which is how the
    # wordmark once shipped with brand and tagline wrapped over two lines,
    # and the title footer with "leaves the machine" hanging under the credit. This
    # estimates the width one line needs and complains at build time instead.
    #
    # It only runs where the caller says oneline=True. Inferring the intent from the
    # box height was tried first and it missed both of the real cases: a 0.3in box at
    # 11pt has room for a second line by arithmetic, which is exactly why nobody
    # noticed the wordmark had taken it. The furniture knows what it is.
    #
    # The advances are per family - mono is 0.6em exactly, the sans averages near
    # 0.54 over mixed-case prose - so the number is an estimate, and it only has to
    # be good enough to catch a box that is out by half an inch.
    ADVANCE = {MONO: 0.60, SANS: 0.54}

    def _check_fits(self, w, paras, size, font):
        for para in paras:
            pieces = para if isinstance(para, list) else [(para, {})]
            line = "".join(text for text, _ in pieces)
            if not line.strip():
                continue
            need = sum(len(text) * self.ADVANCE.get(over.get("font", font), 0.56)
                       * over.get("size", size) / 72 for text, over in pieces)
            if need > w:
                self.warnings.append(
                    f"slide {self.n}: needs {need:.2f}in in a {w:.2f}in box - {line[:58]!r}")

    # -- furniture ----------------------------------------------------------

    def slide_new(self, dark=True, rain=False, credit=True):
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = self._rect(0, 0, W, H, fill=BG if dark else BG2)
        bg.fill.gradient()
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = BG
        stops[0].position = 0.0
        stops[1].color.rgb = BG2
        stops[1].position = 1.0
        bg.fill.gradient_angle = 60.0
        if rain:
            self._rain()
        self.n += 1
        if credit:
            self.watermark()
        return self.slide

    def _rain(self):
        """A faint curtain of glyphs down the right edge. Decoration with a job:
        it makes the deck look like the product instead of a template."""
        for i, col in enumerate(range(9)):
            x = W - 2.6 + col * 0.28
            chars = "\n".join(RAIN[(i * 7 + k * 5) % len(RAIN)] for k in range(11))
            self._text(x, 0.35 + (i % 3) * 0.5, 0.3, H, chars, size=11,
                       color=RGBColor(0x14, 0x3A, 0x22), font=MONO, spacing=1.35, space_after=0)

    def wordmark(self):
        # 4.2in because the line needs 3.1: at 2.0 it broke after "agents" and the
        # running head read as two lines of nothing in particular.
        self._rect(MARGIN, 0.44, 0.1, 0.1, fill=GREEN, shape=MSO_SHAPE.OVAL)
        self._text(MARGIN + 0.2, 0.36, 4.2, 0.3, [[("agentbox", {"bold": True, "color": INK}),
                                                   ("  ·  stop babysitting your agents", {"color": INK3})]],
                   size=11, font=MONO, oneline=True)

    def number(self):
        self._text(W - MARGIN - 1.0, H - 0.62, 1.0, 0.3, f"{self.n:02d}", size=11,
                   color=INK3, font=MONO, align=PP_ALIGN.RIGHT, oneline=True)

    def watermark(self):
        """The author, bottom left, quiet on purpose: a credit the room can read if
        it looks for it, and does not have to read if it does not."""
        self.credits.append(
            self._text(MARGIN, H - 0.62, 4.0, 0.3, AUTHOR, size=10, color=INK3, font=MONO,
                       oneline=True))

    def save(self, path):
        """Write the file, with every credit lifted above the content.

        The credit goes down with the background, so any panel that reaches the bottom
        of a slide is drawn after it and hides it - which is what a licence-required
        credit must not do. Moving the element to the end of its slide's shape tree
        puts it back on top, and costs less than keeping every slide's furniture out of
        that one corner by hand."""
        for tb in self.credits:
            el = tb._element
            parent = el.getparent()
            parent.remove(el)
            parent.append(el)
        self.prs.save(path)

    def kicker(self, text, color=GREEN):
        self._text(MARGIN, 1.18, COLW, 0.3, text.upper(), size=12, color=color, font=MONO,
                   bold=True, spacing=1.0, oneline=True)

    def title(self, text, size=42, y=1.62, color=INK, w=None):
        self._text(MARGIN, y, w or COLW - 1.6, 1.5, text, size=size, color=color, bold=True,
                   spacing=1.06, space_after=0)

    def rule(self, y=2.9, w=1.5, color=GREEN):
        self._rect(MARGIN, y, w, 0.035, fill=color)

    def bullets(self, items, y=3.3, size=19, gap=0.62, w=None, color=INK2):
        """Each item: (bold lead, rest) or a plain string."""
        for i, item in enumerate(items):
            top = y + i * gap
            self._rect(MARGIN + 0.02, top + 0.16, 0.09, 0.09, fill=GREEN, shape=MSO_SHAPE.OVAL)
            if isinstance(item, tuple):
                lead, rest = item
                runs = [[(lead + "  ", {"color": INK, "bold": True}), (rest, {"color": color})]]
            else:
                runs = [[(item, {"color": color})]]
            self._text(MARGIN + 0.34, top, (w or COLW) - 0.34, 0.55, runs, size=size)

    def mono_block(self, lines, x=None, y=4.9, w=None, h=None, label=None):
        x = MARGIN if x is None else x
        w = (COLW if w is None else w)
        h = (0.42 + 0.30 * len(lines)) if h is None else h
        self._rect(x, y, w, h, fill=SURFACE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=EDGE)
        if label:
            self._text(x + 0.28, y + 0.14, w - 0.5, 0.24, label, size=10, color=INK3, font=MONO)
        paras = []
        for ln in lines:
            if ln.startswith("$ "):
                paras.append([("$ ", {"color": GREEN}), (ln[2:], {"color": INK})])
            elif ln.startswith("# "):
                paras.append([(ln, {"color": INK3})])
            elif ln.startswith("→ "):
                paras.append([(ln, {"color": GREEN})])
            else:
                paras.append([(ln, {"color": INK2})])
        self._text(x + 0.28, y + (0.44 if label else 0.22), w - 0.5, h - 0.4, paras,
                   size=13, font=MONO, spacing=1.2, space_after=2)

    def handson(self, what, cmds, note):
        """A demo slide: a badge, what the room is about to see, the command, notes.

        The badge is the only furniture. An earlier version printed the act number of
        the presenter's script next to it, which is a note to the presenter sitting on
        a slide the audience is reading."""
        self.slide_new(rain=True)
        self.wordmark()
        badge = self._rect(MARGIN, 1.1, 2.55, 0.42, fill=GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        badge.adjustments[0] = 0.5
        self._text(MARGIN + 0.26, 1.19, 2.25, 0.3, "LIVE  ·  NOT A MOCK-UP", size=11.5, color=BG,
                   font=MONO, bold=True, oneline=True)
        self.title(what, size=38, y=1.85)
        self.rule(y=3.15)
        self.mono_block(cmds, y=3.62, label="on this machine, now")
        self.notes(note)
        self.number()

    def notes(self, text):
        self.slide.notes_slide.notes_text_frame.text = text.strip()


# ---------------------------------------------------------------- the deck

def build(path):
    d = Deck()

    # 1 ---------------------------------------------------------------- title
    d.slide_new(rain=True)
    d._rect(0, 0, 0.14, H, fill=GREEN)
    d._text(MARGIN, 2.25, 9.4, 1.0, "agentbox", size=104, color=INK, bold=True, spacing=0.95,
            space_after=0)
    d._text(MARGIN + 0.06, 3.62, 9.6, 0.6, "Stop babysitting your agents.", size=30, color=GREEN,
            spacing=1.1, space_after=0)
    d._text(MARGIN + 0.06, 4.42, 9.2, 1.2,
            "A desktop interaction hub for AI agents: a card over whatever you are doing, "
            "a sound that says what kind of thing it is, and an answer that goes straight "
            "back to the code that is blocked on it.", size=18, color=INK2, spacing=1.35)
    # The footer takes the whole column and starts on MARGIN, which is where the
    # author watermark below it starts. Any narrower and 91 characters of mono wrap,
    # which leaves "leaves the machine" hanging over the credit; any further right
    # and the credit reads as indented from a line it is supposed to sit under.
    d._text(MARGIN, H - 1.15, COLW, 0.4,
            [[("one binary  ·  fourteen tools  ·  MCP and CLI  ·  no cloud, no account, "
               "nothing leaves the machine", {"color": INK3})]], size=13, font=MONO,
            oneline=True)
    d.notes("""
This is AgentBox, as in quick question: the thing you type in chat before interrupting a
colleague. It exists so that the agents you run can reach you, and so that being
reached costs you about two seconds.

Everything in the next few minutes is the real software running on one machine. Every
card you are about to see is a real question, and every answer goes straight back to
the code that was waiting for it.
""")

    # 2 ---------------------------------------------------------------- the gap
    d.slide_new()
    d.wordmark()
    d.kicker("the problem")
    d.title("Your agent has been waiting\nforty minutes for one word.")
    d.rule(y=3.35)
    d.bullets([
        ("It stops and waits.", "In a terminal nobody is looking at. Forty minutes of nothing."),
        ("Or it decides for itself.", "And you find out afterwards, in the diff."),
        ("Either way you babysit.", "Which is the opposite of why you started the agent."),
    ], y=3.8, gap=0.78)
    d.notes("""
If you run coding agents, you know this moment. The agent works away for twenty minutes
while you are somewhere else, and then it needs one word from you, and it has nowhere to
put the question.

So it does one of two things, and both cost you. It stops and waits in a terminal nobody
is watching, or it decides for itself and you find out afterwards, in the diff. Either way
you end up sitting next to it, which is the opposite of why you started it.
""")

    # 3 ---------------------------------------------------------------- the idea
    d.slide_new(rain=True)
    d.wordmark()
    d.kicker("the idea")
    d.title("One place every agent\ncan reach you.")
    d.rule(y=3.35)
    d.bullets([
        ("An interruption has a price.", "So the cheap ones cost you nothing but a sound."),
        ("Answering takes about two seconds.", "One keystroke or one click, wherever you were."),
        ("The answer goes back to the code.", "Exit codes and JSON, not a chat transcript."),
    ], y=3.8, gap=0.78)
    d.notes("""
AgentBox is a third option, and a simple one. Every agent on the machine gets one place to
reach you, and being reached is made cheap: a card over whatever you are doing, a sound
that tells you what kind of thing arrived, and an answer that goes straight back to the
code that was waiting for it.

Watch what each interruption costs as we go, because that is the whole design.
""")

    # 4 ------------------------------------------------------- what it buys you
    d.slide_new()
    d.wordmark()
    d.kicker("what it buys you")
    d.title("Hours of unattended work,\nwith the decisions still yours.", size=40)
    d.rule(y=3.35)
    d.bullets([
        ("An agent can work for an hour on its own.", "It asks when a person is genuinely needed."),
        ("You get the decisions, not the steps.", "Approving every command stops being the price."),
        ("Nothing is lost while you are away.", "Every question waits, with a record of what it cost."),
    ], y=3.8, gap=0.78)
    d.notes("""
Here is what that is worth in a working day. An agent that can reach you is an agent you
can leave running, because the moment it needs a person it has somewhere to put the
question instead of stopping.

You stop approving every command it runs, and start answering the few things that actually
need judgement. And whatever arrives while you are away is kept, so coming back to your
desk means reading an inbox rather than reconstructing what happened.
""")

    # 5 ---------------------------------------------------------------- what it is
    d.slide_new()
    d.wordmark()
    d.kicker("what it is")
    d.title("One binary. Fourteen tools.\nEvery agent, every project.")
    d.rule(y=3.3)
    d._text(MARGIN, 3.75, 5.5, 2.4, [
        [("Notify · Ask · Input · Confirm · Veto · Form · Secret · Review",
          {"color": INK, "size": 17})],
        [("Progress · Document · Artifact · Await · Read · Speak", {"color": INK, "size": 17})],
        [("", {})],
        [("MCP for agents, the same fourteen as a CLI for shell scripts, hooks and cron. "
          "One resident daemon owns the windows.", {"color": INK2, "size": 16})],
    ], spacing=1.5)
    d.mono_block([
        "$ claude mcp add --scope user agentbox agentbox mcp",
        "# fourteen tools, in every session, in every project",
        "$ agentbox ask --title 'Deploy?' --option Yes --option Hold",
        "→ Yes            # exit 0, on stdout, ready for a shell if",
    ], x=6.9, y=3.62, w=5.45, label="setup, once")
    d.notes("""
It is one Go binary and a small daemon that stays running. Fourteen tools your agent can
reach over MCP, and the same fourteen from a shell script or a cron job. You register it
once, and from then on every agent in every project can reach you.
""")

    # 6 ------------------------------------------------------- hands-on: earcons
    d.handson("The cheapest interruption is a sound.",
              ["$ agentbox notify --level info    --title 'Dependency scan finished'",
               "$ agentbox notify --level success --title 'Tests are green on main'",
               "$ agentbox notify --level warning --title 'Staging disk is at 86%'",
               "$ agentbox notify --level error   --title 'The nightly build failed'",
               "$ agentbox notify --level urgent  --title 'Production error rate is climbing'"],
              """
Five levels, five sounds. You learn them in a day, and after that you know what kind of
thing happened without stopping what you are doing.

Watch which ones leave on their own. The scan and the green tests took themselves off the
screen, because nothing needs you. The disk warning, the failed build and the production
alert waited until they were read. And the last one arrives even in do-not-disturb, which
is the only reason to have a level above the others.
""")

    # 7 ---------------------------------------------------------- shapes of asking
    d.slide_new()
    d.wordmark()
    d.kicker("the vocabulary")
    d.title("Six shapes of asking,\nand five of them wait for you.", size=38)
    d.rule(y=3.25)
    rows = [
        ("notify", "something happened", "no", "a sound and a card you can ignore"),
        ("ask", "one of two to nine", "yes", "digits 1-9 answer it"),
        ("input", "your words", "yes", "free text, typed and sent"),
        ("confirm", "yes or no", "yes", "exit 0 or exit 1"),
        ("veto", "I am about to…", "yes", "proceeds unless you stop it"),
        ("form", "four answers at once", "yes", "one card, not four interruptions"),
    ]
    y0 = 3.72
    d._text(MARGIN, y0 - 0.34, COLW, 0.26, [[
        ("TOOL              WHAT IT ASKS                BLOCKS    HOW IT IS ANSWERED",
         {"color": INK3})]], size=11, font=MONO)
    for i, (tool, asks, blocks, how) in enumerate(rows):
        y = y0 + i * 0.44
        d._rect(MARGIN, y - 0.04, COLW, 0.38, fill=SURFACE if i % 2 == 0 else None)
        d._text(MARGIN + 0.18, y + 0.03, 2.0, 0.3, tool, size=14, color=GREEN, font=MONO, bold=True)
        d._text(MARGIN + 2.1, y + 0.03, 3.2, 0.3, asks, size=14, color=INK)
        d._text(MARGIN + 5.5, y + 0.03, 1.0, 0.3, blocks, size=14,
                color=AMBER if blocks == "yes" else INK3, font=MONO)
        d._text(MARGIN + 6.9, y + 0.03, 4.4, 0.3, how, size=14, color=INK2)
    d.notes("""
Six shapes of asking, and picking the right one is most of using AgentBox well. Only the
first of them is free - the other five stop the agent until you answer, which is why the
rule AgentBox gives an agent is one line: only ask for a decision you cannot safely make on
your own.
""")

    # 8 -------------------------------------------------- hands-on: two seconds
    d.handson("A decision, in two seconds.",
              ["$ agentbox ask --title 'Where should this release go?' \\",
               "      --option Staging --option Production --option 'Hold it'",
               "→ {\"answered\":true,\"answer\":\"Staging\"}",
               "",
               "$ agentbox input   --title 'What should I tag it?'        # your words",
               "$ agentbox confirm --title 'Run the migration on staging?'  # y or n"],
              """
When an agent genuinely needs you, the only number that matters is how long you take to
answer. About two seconds, and no hunting for the right window.

One click, and the agent already has it. It was blocked; now it is running, and nothing was
left sitting in a terminal.

Sometimes the answer is your words rather than a button, so it gets typed at a keyboard in
the moment it was needed and comes back as a string.

And every card answers from the keyboard too, which matters more than it sounds. If
answering is slower than a key press you will start ignoring the cards, and then none of
this was worth having.
""")

    # 9 ---------------------------------------------------------------- veto
    d.slide_new(rain=True)
    d.wordmark()
    d.kicker("the one you will use most")
    d.title("Silence means yes.\nAnd silence is free.", size=44)
    d.rule(y=3.4)
    d.bullets([
        ("Most of the time the agent knows what to do.", "It just should not do it silently."),
        ("So it says what it is about to do, and counts down.", "You do nothing; it proceeds."),
        ("Forty permission prompts become none.", "And it still never surprises you."),
    ], y=3.85, gap=0.72)
    d.mono_block(["$ agentbox veto --in 10 --title 'Rotating the staging certificate'",
                  "→ proceeding     # exit 0. Stop → 'vetoed', exit 1"], y=5.95, h=0.85)
    d.notes("""
Most of the time the agent already knows what to do. It just should not do it silently.
So it says what it is about to do, counts down, and does it, unless you say no.

Silence means yes, and silence is free. That is what lets an agent work for an hour
without forty permission prompts, and still never surprise you.
""")

    # 10 -------------------------------------------------- hands-on: veto + form
    d.handson("Act unless stopped. Then three answers in one card.",
              ["$ agentbox veto --in 12 --level warning \\",
               "      --title 'Rotating the staging TLS certificate'",
               "# nobody stops it, and that is the point",
               "",
               "$ agentbox form --title 'Release checklist' \\",
               "      --field choice:target:staging,production,canary \\",
               "      --field text:tag --field bool:migrate"],
              """
The certificate expires on Sunday and rotating it is the obvious thing to do, so the agent
says what it is about to do and starts counting. Nobody is going to stop that one, and
nobody had to approve it either. Silence was the answer, and silence was free.

Then the questions it genuinely has get asked together. A target, a version tag, and whether
to run the migration.

Three separate prompts would have been three interruptions; one card is one. An agent that
batches what it needs is an agent you keep installed.
""")

    # 11 ------------------------------------------------------ hands-on: progress
    d.handson("A long job, without a\nnotification per step.",
              ["$ ./reindex.sh | agentbox progress --title 'Reindexing the search corpus'",
               "",
               "# each line of stdin is a percentage and a label:",
               "#   0 reading rows · 45 building the index · 90 verifying · 100 done",
               "",
               "# its own window, in a corner, and it never takes your focus"],
              """
Long jobs used to leave you with two bad options: a notification for every step, or silence
until the end. This is the third one. A bar in its own window, in a corner, that never
takes your focus and never covers the middle of the screen where a question would land.

It reports and you glance. One chime when it finishes, and an error you cannot miss if it
does not.
""")

    # 12 ---------------------------------------------------------------- trust
    d.slide_new()
    d.wordmark()
    d.kicker("the three questions worth asking")
    d.title("Can it leak? Can it phone home?\nCan I switch it off?", size=36)
    d.rule(y=3.25)
    cards = [
        ("A secret never reaches\nthe transcript", GREEN,
         "Masked entry, straight to a 0600 file. The agent gets a path, not a value. "
         "The log records that a secret was asked for - never what it was."),
        ("Agent content cannot\nreach a host", BLUE,
         "An artifact runs with no network at all. An image may name a local file and "
         "nothing else. Both enforced by tests, not by a promise."),
        ("Every switch you\nwould expect", AMBER,
         "Do-not-disturb, per-agent mute, quiet hours - and the one level that pierces "
         "DND can be turned off too."),
    ]
    for i, (head, accent, body) in enumerate(cards):
        x = MARGIN + i * (COLW / 3 + 0.02)
        w = COLW / 3 - 0.28
        d._rect(x, 3.7, w, 2.75, fill=SURFACE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=EDGE)
        d._rect(x, 3.7, w, 0.055, fill=accent)
        d._text(x + 0.3, 4.0, w - 0.6, 0.8, head, size=19, color=INK, bold=True, spacing=1.15)
        d._text(x + 0.3, 4.95, w - 0.6, 1.4, body, size=14, color=INK2, spacing=1.35)
    d.notes("""
These are the three questions anybody serious asks, in this order, and AgentBox has a
straight answer to each. The part worth looking at is the last line of each card: it is
held up by a test in the repository, not by a promise in a readme.
""")

    # 13 ------------------------------------------------- hands-on: secret + dnd
    d.handson("A secret, and the off switch.",
              ["$ agentbox secret --title 'Deploy token for the release' --to-file ./token",
               "→ {\"secret_path\":\"./token\"}     $ stat -c %A ./token  →  -rw-------",
               "",
               "$ agentbox dnd on     # the routine notice never appears; it is in the inbox",
               "$ agentbox notify --level urgent --title 'Production error rate is climbing'",
               "# urgent pierces DND - and you can turn even that off"],
              """
Agents need tokens, and a token must never land in a transcript, a context window or a log.
So the value goes from your keyboard into a file only you can read, and the agent is handed
a path to it. It reads it once and deletes it.

The file is readable by you and nobody else, and AgentBox's own log records that a secret was
asked for. It has never recorded one.

And anything that can interrupt you needs a switch that stops it, or it ends up uninstalled.
That is do-not-disturb: the routine notice you just did not see is in the inbox rather than
on your screen. Held is not lost. One level is still allowed through, and you can turn even
that off.
""")

    # 14 ------------------------------------------------------ hands-on: review
    d.handson("Review the three commands that matter,\nnot all forty.",
              ["$ agentbox review --title 'Approve this change to the API client?' \\",
               "      --diff-file timeout.diff",
               "→ {\"answer\":\"approved\",\"approved\":true}",
               "# 'Request changes' returns rejected, and your note comes back as reply"],
              """
Before an agent changes your code, you get to see the change. Green for what arrives, red
for what goes, and you approve it or send it back with a note in the same two seconds.

This one is a request with no timeout, which pinned a worker for forty minutes on Tuesday.
That is the kind of change worth a person's eye - and it is how you stop watching every
command an agent runs and start reviewing the three that actually matter.
""")

    # 15 ---------------------------------------------------------- the artifact
    d.slide_new(rain=True)
    d.wordmark()
    d.kicker("when words are the wrong shape")
    d.title("When the answer is a number,\na shape, or a selection.", size=38)
    d.rule(y=3.35)
    d.bullets([
        ("Some answers are not a sentence.", "A rollout share, a threshold, a point on a curve."),
        ("So the agent builds the control.", "A small interface in a window, with no network."),
        ("You set it; the agent gets the value.", "One number back, not forty, and it unblocks."),
    ], y=3.8, gap=0.72)
    d.mono_block(["$ agentbox show --artifact rollout.jsx      # the interface it wrote",
                  "→ {\"action\":\"start\",\"percent\":50}        # what you set, back in the code"],
                 y=5.95, h=0.85)
    d.notes("""
Everything so far has been words going back and forth. This one is different. Sometimes the
answer is not a sentence at all: it is how much traffic a new build should take, or where a
threshold belongs.

So the agent writes the control for it, and AgentBox runs it in a window on your desk. You set
it, and the value goes back to the code that was waiting. Forty small moves of a slider
arrive as one number, because forty was never the point.
""")

    # 16 ---------------------------------------------------- hands-on: artifact
    d.handson("An interface the agent wrote,\nfor a decision it cannot make.",
              ["$ id=$(agentbox show --artifact rollout.jsx)",
               "$ agentbox artifact wait --id $id --name rollout      # this blocks",
               "",
               "→ rollout {\"action\":\"start\",\"percent\":50}",
               "# React in a sandbox with no network, and one way out"],
              """
This is the canary decision from a real deploy: how much live traffic the new build should
take. Too little proves nothing, too much and a bad build reaches everybody. Nobody wants to
express that as a sentence in a chat box.

There is an agent parked on this window right now, blocked, waiting for somebody to use it.
The bar underneath the number says what the number actually means: at ten percent that is
four hundred requests a minute, and the whole drag comes back as one value rather than forty.

Half the traffic, and start. That click is what unblocks the deploy, and the agent on the
other side of it now has a number it could never have got from a chat box.
""")

    # 17 ------------------------------------------------ hands-on: document + app
    d.handson("A report worth reading, and the record\nof what everything cost.",
              ["$ agentbox show --title 'Release report' report.md   # tables, maths, charts",
               "",
               "$ agentbox app --tab inbox      # who asked, what for, how long you took",
               "$ agentbox stats --since 24h"],
              """
Agents produce reports, plans, tables and numbers, and in a terminal you scroll past all of
it. This is the report from the release you just approved, in a window built to read it in.

A table of what shipped and what it risked. The arithmetic behind that rollout number, set
as real mathematics. The pipeline it ran, drawn from the agent's own description of it. A
chart of what the week cost in interruptions. And an image read straight off this disk,
because nothing on that page is allowed to touch the network.

Then the record. Every interruption is kept: who asked, what for, and how long you took to
answer, per agent and per project.

That is how you find out which agent is expensive, and turn that one down.

And it is the same window that holds the sounds, the quiet hours and every other switch. A
change there applies while you watch, not at the next restart.
""")

    # 18 ------------------------------------------------------------ the panel
    d.slide_new(rain=True)
    d.wordmark()
    d.kicker("one more thing", color=BLUE)
    d.title("And when you want the agent,\nrather than the other way round.", size=38)
    d.rule(y=3.35, color=BLUE)
    d.bullets([
        ("One hotkey, over whatever you were doing.", "A console rolls down from the top edge."),
        ("It always has a session in it.", "You came with a sentence, not to click New."),
        ("It knows it is inside AgentBox.", "So it reports back instead of waiting in a pipe."),
    ], y=3.8, gap=0.72)
    d.mono_block(["ctrl+alt+`                            # or `agentbox panel show`",
                  "# type into it; a real Claude Code session answers, streaming",
                  "→ plan or full, per session · Esc rolls it back up"], y=5.9, h=0.95)
    d.notes("""
Everything so far has been the agent asking you. This is the same channel in the other
direction: one hotkey, and a console comes down over whatever you were doing, with a live
session already in it.

Two things make it worth using instead of opening a terminal. It is already there, with
no window to find and no directory to change into. And the session knows it is inside
AgentBox, so when you give it something long and walk away, it reaches you when it is done.
""")

    # 19 --------------------------------------------- hands-on: the quake panel
    d.handson("A console, over everything, with the\nagent already in it.",
              ["ctrl+alt+`         # it rolls down, wherever you were",
               "",
               "> In one sentence: what is agentbox, and when should you interrupt me?",
               "",
               "# a real Claude Code session answers, streaming, and the same key",
               "# puts it away. Plan or full, per session."],
              """
Watch the top of the screen. One key, and there it is, centred on the monitor, half its
height, with a session waiting.

That is a real Claude Code session on this machine, and the reply is streaming as it is
written.

Notice what it says about interrupting me. It was told at startup that it is running
inside AgentBox, so it treats a question as something to raise here, and a finished job as
something worth telling you about. And the same key puts it away.
""")

    # 20 -------------------------------------------------------- what is next
    d.slide_new()
    d.wordmark()
    d.kicker("what is next", color=BLUE)
    d.title("It is in daily use here,\nand it is still being built.", size=40)
    d.rule(y=3.3, color=BLUE)
    nexts = [
        ("Away delivery", BLUE,
         "A relay, so a question that matters can reach you when you are not at this "
         "desk at all - the one thing a local-only tool cannot do today."),
        ("Sessions without a terminal", GREEN,
         "Start, watch and answer agent sessions from the panel itself, with the progress "
         "of each one visible while it runs."),
        ("Wayland, and packages", AMBER,
         "The window work that X11 gets today, on Wayland - and distribution packages, so "
         "installing it is one command from a repository."),
    ]
    for i, (head, accent, body) in enumerate(nexts):
        x = MARGIN + i * (COLW / 3 + 0.02)
        w = COLW / 3 - 0.28
        d._rect(x, 3.78, w, 2.5, fill=SURFACE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=EDGE)
        d._rect(x, 3.78, w, 0.055, fill=accent)
        d._text(x + 0.3, 4.08, w - 0.6, 0.5, head, size=19, color=INK, bold=True, spacing=1.15)
        d._text(x + 0.3, 4.72, w - 0.6, 1.4, body, size=14, color=INK2, spacing=1.35)
    d._text(MARGIN, 6.5, COLW, 0.4, [[
        ("Multi-select questions, a preview of what is queued, richer forms: the vocabulary "
         "keeps growing.", {"color": INK3})]], size=14)
    d.notes("""
Two honest things about where this is. It is in daily use on this machine, and it is still
being built, so what you have seen is a working tool rather than a finished product.

Three things are coming that change what it can do. Away delivery, so a question that matters
can reach you when you are not at this desk at all. Running whole agent sessions from that
panel, not just talking to one. And the same window handling on Wayland, with packages so
installing it is one command.

Underneath that, the vocabulary keeps growing: more shapes of asking, and a preview of what
is queued behind the card you are looking at.
""")

    # 21 ---------------------------------------------------------------- get it
    d.slide_new()
    d.wordmark()
    d.kicker("how you get it")
    d.title("Four commands, then forget\nit is there.", size=40)
    d.rule(y=3.3)
    d.mono_block([
        "$ make bootstrap    # toolchain, GTK4/WebKit headers, voice, config",
        "$ make install      # binary, launcher, systemd user service",
        "$ systemctl --user enable --now agentbox.service",
        "$ claude mcp add --scope user agentbox agentbox mcp",
        "",
        "# and for an agent, the whole manual:",
        "$ agentbox docs agent",
    ], y=3.75, label="a machine that has never seen it")
    d._text(MARGIN, 6.5, COLW, 0.5, [[
        ("make doctor", {"font": MONO, "color": GREEN}),
        ("  tells you what is missing, and installs nothing.", {"color": INK2})]], size=15)
    d.notes("""
Getting it is four lines. Bootstrap installs what a fresh machine needs, including a
speech engine and a voice. Install puts the binary where your desktop can find it. One
more line registers it with Claude Code, for every project at once. After that you forget
it is there, until it has a question.
""")

    # 22 ---------------------------------------------------------------- close
    d.slide_new(rain=True)
    d._rect(0, 0, 0.14, H, fill=GREEN)
    d._text(MARGIN, 2.35, 10.4, 2.0,
            "Your agents run unattended.\nYou get pulled in for the decisions\nonly you can make.",
            size=38, color=INK, bold=True, spacing=1.15, space_after=0)
    d.rule(y=4.9)
    d._text(MARGIN, 5.25, 10.0, 0.8,
            "agentbox · one binary, a daemon, and fourteen ways to reach you. "
            "No cloud, no account, nothing leaving this machine.",
            size=17, color=INK2, spacing=1.35)
    d._text(MARGIN, 6.35, 10.0, 0.4, [[("agentbox docs agent  ·  the manual an agent reads "
                                        "before it interrupts you for the first time.",
                                        {"color": INK3, "italic": True})]], size=14)
    d.notes("""
That was a release going out, a job being watched, a secret being handed over and a diff
being approved - about fifteen interruptions, and not one of them took more than a couple
of seconds to deal with.

That is the whole idea. Your agents run unattended, and you get pulled in for the decisions
that are actually yours. One binary, on your own machine, and nothing leaves it. Thanks for
watching.
""")

    d.save(path)
    return path, d.warnings


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("-o", "--out", default=os.path.join("docs", "agentbox-showcase.pptx"))
    args = ap.parse_args()
    out, warnings = build(args.out)
    print(f"wrote {out} ({os.path.getsize(out) / 1024:.0f} kB)")
    # Loud on purpose. A wrapped label is not visible to the person who wrote the
    # slide and is the first thing an audience sees.
    for w in warnings:
        print(f"  wraps: {w}")
    if warnings:
        print(f"{len(warnings)} single-line box(es) do not fit; widen the box or shorten the text")
