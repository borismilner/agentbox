#!/usr/bin/env python3
"""Perform the narrated showcase, in sync with the deck.

    python3 tools/showcase/perform.py --dry-run
    python3 tools/showcase/perform.py                 # the whole thing
    python3 tools/showcase/perform.py --from 7 --to 9 # a rehearsal

The narration is not in this file. It lives in the deck's speaker notes, which
tools/showcase/deck.py generates, so the words on screen and the words in the air
are one artifact and cannot drift. A blank line in a note is a beat boundary: this
reads the beats in order and speaks them with `agentbox say --wait`, which returns when
the audio has actually stopped rather than after a guess. That is what keeps the
narration synchronised with the slides for eight minutes without a single sleep
tuned by hand.

What this file holds is the other half: what happens on screen between the beats.
Each slide's entry is a list of steps, run in order, and "beat" means "say the next
sentence of this slide's narration". The commands are the ones in docs/showcase.md,
which is the plan a human reads; this is the same sequence, executed.

Slides advance with the right arrow at the end of each entry, so the deck follows the
voice rather than a timer. Start it with the slideshow already fullscreen (see
tools/showcase/record.sh, which prepares the desk and rolls the camera).

The synthetic input below is stagecraft and never the subject: it is how the cards get
answered with nobody in the frame, and the deck does not mention it. Do not put it back
into a slide or a narration beat.
"""

import argparse
import os
import shlex
import subprocess
import sys
import time

AGENTBOX = os.path.expanduser("~/.local/bin/agentbox")
RECORD = os.path.join("tools", "showcase", "record.sh")
DECK = os.path.join("docs", "agentbox-showcase.pptx")

# Who the cards come from. The tour is one team's afternoon rather than a feature
# list, so four different agents ask over the course of it: the inbox and the
# insights tab then show a real distribution instead of one demo name, which is the
# whole argument of that slide.
PROJECT = "checkout-api"
BOT = "release-bot"

# The pointer's parking spot, computed from the recorded monitor at run time. A
# pointer left in the middle of a slide covers the words it is there to introduce,
# and one left over the presenter's own control strip makes it appear.
PARK = None


# ---------------------------------------------------------------- the steps

def sh(cmd, background=False, check=False):
    """Run a shell line. Background is for the blocking card commands, which do not
    return until somebody answers them - that somebody being a later step.

    shell=True is deliberate and safe here: every command is a literal in the table
    below, and several of them need a shell (a trailing &, a pipe, a redirect). No
    caller-supplied string reaches this function - the only command-line argument this
    script takes is a pair of integers for the pointer.

    A command ending in `&` is run without pipes. capture_output hands the child a
    pipe that its own background grandchild inherits, so run() waits for that pipe to
    close - which means it waits for the "backgrounded" command to finish. A veto
    written as ("run", "agentbox veto --in 12 ... &") therefore blocked for its whole
    twelve seconds, the card came and went inside that one step, and the narration
    that followed talked about a card that was no longer there. Seen on camera."""
    if background or cmd.rstrip().endswith("&"):
        return subprocess.Popen(cmd, shell=True, stdout=subprocess.DEVNULL,
                                stderr=subprocess.DEVNULL)
    return subprocess.run(cmd, shell=True, capture_output=True, text=True, check=check)


def drive(script):
    """Feed a drive script to the daemon, which owns the synthetic input."""
    return subprocess.run([AGENTBOX, "drive", "run", "-"], input=script, text=True,
                          capture_output=True)


# ------------------------------------------------------------------ the stage
#
# Three things about the machine have to hold for a take to be usable, and all
# three failed silently in the recording of 2026-07-25 22:33. They are checked here
# instead of being hoped for.

STATE = os.path.join(os.environ.get("XDG_RUNTIME_DIR", "/tmp"), "agentbox-showcase-record")


def stage_region():
    """The monitor being recorded, as (w, h, x, y), or None outside a take."""
    try:
        parts = open(os.path.join(STATE, "region")).read().split()
        return int(parts[1]), int(parts[2]), int(parts[3]), int(parts[4])
    except Exception:
        return None


def _geometry(win):
    r = subprocess.run(["xwininfo", "-id", win], capture_output=True, text=True)
    g = {}
    for line in r.stdout.splitlines():
        if "Absolute upper-left X" in line:
            g["x"] = int(line.split()[-1])
        elif "Absolute upper-left Y" in line:
            g["y"] = int(line.split()[-1])
        elif line.strip().startswith("Width:"):
            g["w"] = int(line.split()[-1])
        elif line.strip().startswith("Height:"):
            g["h"] = int(line.split()[-1])
    return g


def stage_window():
    """The fullscreen window covering the recorded monitor, as an X id, or None.
    Outside a take there is no region file, so any fullscreen window counts."""
    reg = stage_region()
    r = subprocess.run(["wmctrl", "-l"], capture_output=True, text=True)
    for line in r.stdout.splitlines():
        win = line.split()[0] if line.split() else ""
        if not win:
            continue
        st = subprocess.run(["xprop", "-id", win, "_NET_WM_STATE"],
                            capture_output=True, text=True)
        if "FULLSCREEN" not in st.stdout:
            continue
        if reg is None:
            return _xid(win)
        g = _geometry(win)
        if (g.get("w"), g.get("h"), g.get("x"), g.get("y")) == reg:
            return _xid(win)
    return None


def deck_is_fullscreen():
    """True while a fullscreen window covers the monitor being recorded.

    The slideshow is its own window, not the editor's, and it leaves fullscreen
    without telling anybody. In the 22:33 take it left at 13:04 and the last four
    minutes of the film are the editor with the deck's last slide in it: nine wheel
    notches meant for a document window landed on the slideshow, each one advanced a
    slide, and the show ran off its end and closed itself. Nothing noticed for four
    minutes, which is the actual defect - the wheel was only the trigger."""
    return stage_window() is not None


def _xid(s):
    try:
        return int(s, 16)
    except ValueError:
        return None


def stacking_order():
    """Every mapped window as an X id, bottom to top, from the root window's
    _NET_CLIENT_LIST_STACKING - the one list that says who is in front of whom."""
    r = subprocess.run(["xprop", "-root", "_NET_CLIENT_LIST_STACKING"],
                       capture_output=True, text=True)
    if "#" not in r.stdout:
        return []
    ids = (_xid(t.strip()) for t in r.stdout.split("#", 1)[1].split(","))
    return [i for i in ids if i is not None]


def window_id(title):
    """The X id of the window with exactly this title, or None. wmctrl and xprop
    write the same id with different zero padding, so ids compare as numbers."""
    want = title.lstrip("=")
    r = subprocess.run(["wmctrl", "-l"], capture_output=True, text=True)
    for line in r.stdout.splitlines():
        parts = line.split(None, 3)
        if len(parts) == 4 and parts[3].strip() == want:
            return _xid(parts[0])
    return None


def above_stage(title, seconds=6.0):
    """None once the window stands in front of the stage, else what is wrong.

    "Appeared" has to mean visible, and nothing else here can tell. agentbox drive where
    asks the daemon, which knows the window it made and not whether anyone can see
    it; the progress bar the uploaded take is missing existed the whole of slide 11 -
    mapped, alive, and under the fullscreen slideshow. Only the camera knew, and the
    camera reported it days later. The stacking list knows at the moment it matters."""
    deadline = time.time() + seconds
    verdict = "never appeared"
    while time.time() < deadline:
        win = window_id(title)
        order = stacking_order()
        if win is not None and win in order:
            stage = stage_window()
            if stage is None or stage not in order or order.index(win) > order.index(stage):
                return None
            verdict = "on screen but UNDER the fullscreen stage"
        time.sleep(0.3)
    return verdict


# There is no input-source handling here any more. agentbox locks the planned keyboard
# group around every synthetic key press itself (internal/hand/xkb.go), and the
# gsettings write this file used to make was measured doing nothing: GNOME 46
# ignores the deprecated `input-sources current` key in both directions.

MAX_SPOKEN = 200  # AgentBox's own [speech] max_chars is 240; stay clear of it


def sentences(text):
    """Split a beat into speakable sentences.

    AgentBox truncates a spoken line at max_chars - 240 by default - because a spoken
    line is meant to be a sentence and not a document. A 360-character beat handed
    over whole comes out with its last sentence missing, which is how the first
    rehearsal lost "it stops and waits in a terminal nobody is watching". So the beat
    is broken on sentence ends here, where the split can be clean, rather than by a
    rune count somewhere downstream.

    Splitting also buys the pacing: each sentence is its own `say --wait`, so the
    pause between two sentences is a real pause and the slide can turn on one.
    """
    out, cur = [], ""
    for word in text.split():
        cur = f"{cur} {word}".strip()
        ends = word.endswith((".", "!", "?")) and not word.endswith(("...",))
        if ends and len(cur) > 40 or len(cur) > MAX_SPOKEN:
            out.append(cur)
            cur = ""
    if cur:
        out.append(cur)
    return out


def say(line):
    """Speak one line and return when it has been heard. The whole synchronisation of
    this performance is this one call: no sleep, no words-per-second guess."""
    return subprocess.run([AGENTBOX, "say", "--wait", "--timeout", "180", line],
                          capture_output=True, text=True)


def speak_beat(text, dry, tag):
    for line in sentences(text):
        print(f"  [{tag}] say: {line[:78]}")
        if not dry:
            say(line)


def card_geometry(title="=agentbox", tries=40):
    """Where the current card is, or None. agentbox drive where asks the daemon rather
    than the window manager, because wmctrl reports doubled coordinates here."""
    for _ in range(tries):
        r = subprocess.run([AGENTBOX, "drive", "where", title], capture_output=True, text=True)
        if r.returncode == 0 and len(r.stdout.split()) == 4:
            return [int(v) for v in r.stdout.split()]
        time.sleep(0.25)
    return None


def submit(script):
    """Answer a card, and check that it actually got answered.

    A measured pixel offset is only true for the card it was measured on. The Send
    button of the secret card was clicked at the form card's Submit offset for one
    take, missed by about ten pixels, and the card then sat on screen until it
    expired - ninety seconds of a recording spent watching a stuck dialog.

    So: click, then look. If the card is still there, press Return, which every card
    accepts while a field has focus (Card.svelte's key handler routes it to the same
    send() the button calls). The click is what the audience should see; Return is
    what makes sure the demo moves on.

    The looking has to outlast the FR28 undo grace. An answered card stays mapped for
    grace_ms - 3000 by default - before the answer ships, so a 0.7s check calls every
    successful click a miss and sends a Return into a card that is already answered.
    This polls until the window really goes, and only then decides it missed.
    """
    drive(script)
    if wait_gone(4.5):
        return "click"
    drive("key Return\n")
    if wait_gone(4.5):
        return "click missed, Return answered it"
    return "STILL ON SCREEN"


def wait_gone(seconds):
    """True once no card is on screen, or False after waiting that long."""
    deadline = time.time() + seconds
    while time.time() < deadline:
        if card_geometry(tries=1) is None:
            return True
        time.sleep(0.25)
    return False


def window_gone(title, seconds=2.5):
    """True once no window matches title, or False after waiting that long."""
    deadline = time.time() + seconds
    while time.time() < deadline:
        r = subprocess.run([AGENTBOX, "drive", "where", title], capture_output=True, text=True)
        if r.returncode != 0:
            return True
        time.sleep(0.25)
    return False


def close_window(title, script):
    """Close a viewer, artifact or app window, and make sure it actually went.

    The click is what the audience should see. The check is because a click that
    misses leaves the window standing in front of everything for the rest of the film:
    slide 17's report window survived its own close click in the 22:33 take and was
    still on screen half an hour later. wmctrl's WM_DELETE_WINDOW is the fallback - it
    needs no focus and no coordinates, so it works when a pixel offset does not."""
    drive(script)
    if window_gone(title):
        return "click"
    ids = subprocess.run(["wmctrl", "-l"], capture_output=True, text=True).stdout
    want = title.lstrip("=")
    for line in ids.splitlines():
        parts = line.split(None, 3)
        if len(parts) == 4 and parts[3].strip() == want:
            subprocess.run(["wmctrl", "-i", "-c", parts[0]], capture_output=True)
            break
    if window_gone(title):
        return "click missed, closed by the window manager"
    return "STILL ON SCREEN"


def clear_cards():
    """Dismiss anything still on screen. A card left over from a missed click would
    sit on top of the next act, so every act that opens one also closes it."""
    for _ in range(6):
        if card_geometry(tries=1) is None:
            return True
        sh(f"{AGENTBOX} summon")
        time.sleep(0.4)
        drive("key shift+Escape\n")
        time.sleep(0.5)
    return card_geometry(tries=1) is None


def park():
    if PARK:
        drive(f"screen\nmove {PARK[0]} {PARK[1]}\n")


# ---------------------------------------------------------------- the plan
#
# Step forms:
#   "beat"              say this slide's next narration beat, and wait for it
#   ("run", cmd)        run a command and wait for it
#   ("bg", cmd)         start a blocking command in the background
#   ("drive", script)   synthetic input
#   ("submit", script)  synthetic input that must answer a card, verified
#   ("close", title, script)  input that must close that window, verified
#   ("wait_card",)      wait until a card is actually on screen
#   ("above", title)    wait until that window is on screen AND in front of the stage
#   ("hold", seconds)   let the audience read
#   ("clear",)          dismiss whatever is on screen
#   ("park",)           pointer out of the way
#
# The pointer parks itself after a run of input steps, never between two of them: a
# cursor that flies to the corner after every field of a form is filled looks like a
# glitch, where one that fills the form and then steps out of the way looks like a
# hand. run_slide decides that by looking at whether the next step is input too.

def ident(agent=BOT):
    return f"--agent {agent} --project {PROJECT}"


def notify(level, title, body, speak, agent=BOT):
    return ("run", f"{AGENTBOX} notify --level {level} --title {shlex.quote(title)} "
                   f"--body {shlex.quote(body)} --speak {shlex.quote(speak)} {ident(agent)}")


# The diff of the change the review card asks about: a real defect with a real
# consequence, small enough to read from the back of a room.
#
# Keep every line under about 56 visible columns, tabs included. The card is 470px
# wide and scrolls a longer line sideways instead of wrapping it, so the first draft
# of this diff - with full internal/api/ paths - showed the audience "b/internal/api/
# client." with the rest off the edge.
DIFF = r"""printf '%s\n' \
 'diff --git a/api/client.go b/api/client.go' \
 '--- a/api/client.go' \
 '+++ b/api/client.go' \
 '@@ -18,7 +18,10 @@ func New(base string) *Client {' \
 ' 	return &Client{' \
 ' 		base: base,' \
 '-		http: &http.Client{},' \
 '+		// No deadline pins a worker until restart.' \
 '+		http: &http.Client{Timeout: 10 * time.Second},' \
 ' 	}' \
 ' }' > /tmp/agentbox-showcase-timeout.diff"""

# The progress bar's source. Seventeen updates at 1.6s is about half a minute, which
# is what the two narration beats over it take - a bar that finishes early leaves the
# slide dead, and one that is still crawling when the voice stops looks stuck.
REINDEX = (
    "for p in 0 4 9 15 22 28 35 41 48 55 62 70 77 84 91 96 100; do "
    "  case $p in "
    "    0)   echo \"$p reading rows from orders\" ;; "
    "    22)  echo \"$p building the index\" ;; "
    "    55)  echo \"$p building the index (2 of 3 shards)\" ;; "
    "    84)  echo \"$p verifying against the old index\" ;; "
    "    100) echo \"$p done: 12,400 rows\" ;; "
    "    *)   echo \"$p working\" ;; "
    "  esac; sleep 1.6; done | "
    f"{AGENTBOX} progress --title 'Reindexing the search corpus' {ident()}"
)


SLIDES = {
    # 1 - the title. Nothing but the promise, spoken over a still slide.
    1: ["beat", "beat"],

    # 2 - the problem, with the tour's own first card under it. It says what is about
    # to happen, because an audience that knows the count watches for the cost.
    2: [notify("info", "This is AgentBox",
               "An agent is about to interrupt you about fifteen times, on purpose.",
               "This is AgentBox, and it is about to interrupt you about fifteen times, on purpose."),
        ("wait_card",), ("hold", 1.5), "beat", ("clear",), "beat"],

    # 3 - the idea.
    3: ["beat", "beat"],

    # 4 - what it buys you.
    4: ["beat", "beat"],

    # 5 - the shape of it.
    5: ["beat"],

    # 6 - the five levels, one at a time, from four different agents. Each one has to
    # land alone or the sound design is a blur, so they are held rather than fired in a
    # burst; each card speaks its own line, so the holds are never silent.
    6: ["beat",
        notify("info", "Dependency scan finished", "214 packages, no new advisories.",
               "Dependency scan finished. Nothing new.", "dependency-bot"),
        ("hold", 3.5), ("clear",),
        notify("success", "Tests are green on main", "395 tests in 4 minutes 12 seconds.",
               "All tests green on main.", "test-runner"),
        ("hold", 3.5), ("clear",),
        notify("warning", "Staging disk is at 86%",
               "Worth a look this week. It waits until you read it.",
               "Staging disk is at eighty six percent.", "oncall-helper"),
        ("hold", 3.5), ("clear",),
        notify("error", "The nightly build failed",
               "The Postgres client library is missing on the runner.",
               "The nightly build failed.", "test-runner"),
        ("hold", 3.5), ("clear",),
        notify("urgent", "Production error rate is climbing",
               "5xx up twelvefold in four minutes. Urgent pierces do-not-disturb.",
               "Production error rate is climbing. This one is urgent.", "oncall-helper"),
        ("hold", 4.5), ("clear",), "beat"],

    # 7 - the vocabulary table. Words only.
    7: ["beat"],

    # 8 - the three ways a question comes back: a button, a line of text, one key.
    # The click offsets are the measured ones from docs/showcase.md, relative to the
    # card window, and they only hold for a one-line body.
    8: ["beat",
        ("bg", f"{AGENTBOX} ask --title 'Where should this release go?' "
               f"--body 'Tests are green and the changelog is written.' "
               f"--option Staging --option Production --option 'Hold it' "
               f"--speak 'Where should this release go?' "
               f"{ident()} --timeout 120 --json"),
        ("wait_card",), ("hold", 2.5),
        # -74, not -94: this body is one line, which makes the card 214px rather than
        # 234, and the option chips sit a fixed distance above the footer. Measured.
        ("submit", "window =agentbox\nspeed 1.15\nclick 60 -74\n"),
        ("hold", 1.0), "beat", ("clear",),
        ("bg", f"{AGENTBOX} input --title 'What should I tag it?' "
               f"--body 'Free text. Whatever you type comes back as a string.' "
               f"{ident()} --timeout 120 --json"),
        ("wait_card",), ("hold", 1.6),
        ("drive", "window =agentbox\nclick 50% -123\nwait 200\nwpm 250\ntype 2026.7.3\n"),
        ("submit", "key Return\n"), ("hold", 0.8), "beat",
        ("bg", f"{AGENTBOX} confirm --title 'Run the migration on staging?' "
               f"--body 'Two columns, no data loss, and it is reversible.' "
               f"{ident()} --timeout 120 --json"),
        ("wait_card",), ("hold", 2.0), "beat",
        ("run", f"{AGENTBOX} summon"), ("submit", "key y\n"),
        ("hold", 1.0), ("clear",)],

    # 9 - act unless stopped: the one a viewer will use every day.
    9: ["beat", "beat"],

    # 10 - the countdown runs out under the narration, then the three questions the
    # agent does have arrive as one card.
    # A veto goes through ("bg", ...): it prints `proceeding` when it elapses, and
    # ("run", "... &") would hold this step for the whole countdown - see sh().
    10: [("bg", f"{AGENTBOX} veto --in 12 --level warning "
                 f"--title 'Rotating the staging TLS certificate' "
                 f"--body 'The old one expires on Sunday. Say the word and I hold.' "
                 f"--speak 'Rotating the staging certificate in twelve seconds, unless you "
                 f"stop me.' {ident()}"),
         ("wait_card",), ("hold", 1.5), "beat", ("clear",),
         ("bg", f"{AGENTBOX} form --title 'Release checklist' "
                f"--field choice:target:staging,production,canary "
                f"--field text:ticket --field bool:migrate "
                f"--speak 'Three things before I ship it.' "
                f"{ident()} --timeout 120 --json"),
         ("wait_card",), ("hold", 1.5), "beat",
         ("drive", "window =agentbox\nclick 64% -212\nwait 300\nclick 64% -123\n"),
         ("drive", "window =agentbox\nclick 64% -162\nwait 200\ntype REL-482\n"),
         ("drive", "window =agentbox\nwait 200\nclick 33 -126\n"),
         ("submit", "window =agentbox\nwait 300\nclick -52 -92\n"),
         ("hold", 1.0), "beat", ("clear",)],

    # 11 - the progress bar: its own window, its own corner, and it never takes focus.
    # The above step is why the uploaded take has to be redone: the bar existed for
    # this whole slide and the film never saw it, because it sat under the slideshow.
    11: [("bg", REINDEX), ("above", "agentbox · progress"), ("hold", 2.5),
         "beat", "beat", ("hold", 2.0)],

    # 12 - the three questions anybody serious asks.
    12: ["beat"],

    # 13 - a secret that never reaches the transcript, then the off switch.
    13: [("bg", f"{AGENTBOX} secret --title 'Deploy token for the release' "
                f"--body 'Masked as you type. The value never travels back over the wire.' "
                f"--speak 'The deploy token, please. It goes into a file, not into my "
                f"context.' --to-file /tmp/agentbox-showcase-token {ident()} "
                f"--timeout 120 --json"),
         ("wait_card",), ("hold", 1.5),
         ("drive", "window =agentbox\nclick 50% -150\nwait 200\n"
                   "wpm 260\ntype ghp_8f3c1d9a24b7e05f1c\n"),
         "beat",
         # Return, not a click: this card's Send sits at its own offset, and the one
         # borrowed from the form card missed it by ten pixels in an earlier take.
         # It also shows the keyboard answering, which the narration keeps promising.
         ("submit", "key Return\n"), ("hold", 0.8), "beat",
         ("run", "stat -c '%A  %n' /tmp/agentbox-showcase-token; rm -f /tmp/agentbox-showcase-token"),
         ("clear",),
         ("run", f"{AGENTBOX} dnd on"),
         notify("info", "Changelog drafted for 2026.7.3",
                "Held. It is in the inbox rather than on your screen.", ""),
         ("hold", 2.0), "beat",
         notify("urgent", "Production error rate is climbing",
                "Urgent pierces do-not-disturb, unless you turn even that off.",
                "This one came through anyway, because urgent is allowed to.", "oncall-helper"),
         ("wait_card",), ("hold", 3.5), ("clear",), ("run", f"{AGENTBOX} dnd off"),
         # Turning do-not-disturb off *delivers* what it held: the changelog notice
         # from the top of this slide arrives the moment the switch flips. Left alone
         # it drifts into the next slide and sits top-right over the review card - and
         # in the takes before toasts had their own title ("agentbox · toast"), it also made
         # slide 14's answered card look like it was still on screen. Take it here, on
         # the slide whose subject it is.
         ("hold", 1.5), ("clear",)],

    # 14 - the review card. Both beats are spoken with the diff on screen, which is
    # also how long a person needs to actually read it.
    14: [("run", DIFF),
         ("bg", f"{AGENTBOX} review --title 'Approve this change to the API client?' "
                f"--body 'A request with no timeout pinned a worker for forty minutes.' "
                f"--diff-file /tmp/agentbox-showcase-timeout.diff "
                f"--speak 'One change to review before I ship it.' "
                f"{ident()} --timeout 180 --json"),
         ("wait_card",), ("hold", 2.0), "beat", "beat",
         ("submit", "window =agentbox\nclick -55 -59\n"),  # Approve
         ("hold", 1.2), ("clear",)],

    # 15 - when the answer is not words.
    15: ["beat", "beat"],

    # 16 - the artifact: an interface the agent wrote, with an agent blocked on it.
    # The window is WebKit and takes a moment to paint, so the first hold is long
    # enough that the narration never talks over a blank frame.
    16: [("run", f"{AGENTBOX} show --artifact --title 'Canary rollout' "
                 f"tools/showcase/console.jsx > /tmp/agentbox-showcase-art"),
         ("bg", "id=$(tail -1 /tmp/agentbox-showcase-art); "
                f"{AGENTBOX} artifact wait --id \"$id\" --name rollout --timeout 180 "
                "> /tmp/agentbox-showcase-rollout"),
         ("above", "agentbox · Canary rollout"), ("hold", 4.0), "beat",
         # The artifact and the viewer are titled "agentbox · <what>", not "agentbox" - only
         # the card is bare now, so a =agentbox selector finds neither.
         # The track spans 36..864 in a 900-wide window, so 50% of the window is
         # exactly half of the range: the number the narration says out loud.
         ("drive", "window =agentbox · Canary rollout\ndrag 14% 47% 50% 47%\n"),
         ("hold", 1.2), "beat",
         ("drive", "window =agentbox · Canary rollout\nclick 26% 86%\n"),
         ("hold", 1.5), "beat",
         ("run", "cat /tmp/agentbox-showcase-rollout"),
         ("close", "=agentbox · Canary rollout",
          "window =agentbox · Canary rollout\nclick -19 18\n")],

    # 17 - the report the agent wrote, and then the record of what everything cost.
    # Three scrolls of nine take the report from its first heading to its last line,
    # measured on this content: the image at the bottom is named in the narration, so
    # it has to actually be reached.
    17: [("run", f"{AGENTBOX} show --title 'Release report' tools/showcase/tour.md"),
         ("above", "agentbox · Release report"), ("hold", 4.0), "beat",
         # Every scroll moves into the window first. `window T` only sets the frame
         # that coordinates are read in - it does not move the pointer - and a wheel
         # notch goes wherever the pointer happens to be. The pointer was parked on
         # the slideshow by the step before, so nine notches meant for this document
         # advanced nine slides and closed the show. lint_plan() now refuses a scroll
         # that has no move in front of it.
         ("drive", "window =agentbox · Release report\nmove center 45%\nscroll 9\n"),
         ("hold", 1.0), "beat",
         ("drive", "window =agentbox · Release report\nmove center 45%\nscroll 9\n"),
         ("hold", 2.2),
         ("drive", "window =agentbox · Release report\nmove center 45%\nscroll 9\n"),
         ("hold", 2.2),
         ("close", "=agentbox · Release report",
          "window =agentbox · Release report\nclick -19 18\n"),
         ("run", f"{AGENTBOX} app --tab inbox"), ("above", "agentbox · app"), ("hold", 3.5), "beat",
         # The rail is session, inbox, history, document, with settings at the bottom.
         # History is the one that answers "what did this cost me"; the document
         # surface would open an empty reader, so the walk skips it.
         ("drive", "window =agentbox · app\nclick 27 141\n"), ("hold", 3.0), "beat",   # history
         ("drive", "window =agentbox · app\nclick 27 -29\n"), ("hold", 2.0), "beat",   # settings
         ("close", "=agentbox · app", "window =agentbox · app\nclick -21 18\n")],

    # 18 - the other direction: you reach the agent.
    18: ["beat", "beat"],

    # 19 - the panel, live, with a real session in it.
    19: [("drive", "key ctrl+alt+grave\n"), ("above", "agentbox · panel"),
         ("hold", 2.0), "beat",
         ("drive", "window =agentbox · panel\nclick 50% -96\n"),
         ("drive", "window =agentbox · panel\nwpm 420\n"
                   "type In one sentence: what is agentbox, and when should you interrupt me\n"),
         ("drive", "key Return\n"), ("hold", 1.0), "beat", ("hold", 6.0), "beat",
         ("drive", "key ctrl+alt+grave\n")],

    # 20 - what is next.
    20: ["beat", "beat", "beat"],

    # 21 - four lines to install it.
    21: ["beat"],

    # 22 - the close, and the last card finishes the story the tour told.
    22: ["beat",
         notify("success", "Release 2026.7.3 is live",
                "Canary at 50% for thirty minutes, error rate flat. Nothing needs you.",
                "Release twenty twenty six point seven point three is live, and nothing "
                "needs you."),
         ("wait_card",), ("hold", 2.5), "beat", ("clear",)],
}


# ---------------------------------------------------------------- narration

def beats_from_deck(path=DECK):
    """The narration, per slide, split on blank lines. Read from the deck so the
    spoken words and the slide are the same artifact."""
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("python-pptx is missing: use .venv-deck/bin/python (make deck creates it)")
    prs = Presentation(path)
    out = {}
    for i, slide in enumerate(prs.slides, 1):
        text = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        # One line per beat, with the source's line breaks folded away: the engine
        # takes one line per utterance.
        out[i] = [" ".join(b.split()) for b in text.split("\n\n") if b.strip()]
    return out


# ---------------------------------------------------------------- the lint
#
# Two mistakes in a drive script are invisible until they are on film, so they are
# refused here instead. Both cost a take.

COORD_OPS = ("click", "double", "drag", "move")


def lint_plan():
    """Complaints about the plan itself, as a list of strings."""
    bad = []
    for n, steps in sorted(SLIDES.items()):
        for step in steps:
            if not (isinstance(step, tuple) and step[0] in INPUT_STEPS):
                continue
            script = step[2] if step[0] == "close" else step[1]
            ops = [ln.split()[0] for ln in script.splitlines() if ln.strip()]
            # A wheel notch is delivered wherever the pointer is standing. `window T`
            # sets the coordinate frame and moves nothing, so a scroll with no move in
            # front of it scrolls whatever the last step left the pointer over - which
            # on this stage is the slideshow, where a notch means "next slide".
            if "scroll" in ops and "move" not in ops:
                bad.append(f"slide {n}: a scroll with no move before it - the wheel "
                           f"would land wherever the pointer is: {script.split(chr(10))[0]!r}")
            # A click with no window and no `screen` is read in whatever frame the
            # previous script left behind.
            if any(o in COORD_OPS for o in ops) and not ("window" in ops or "screen" in ops):
                bad.append(f"slide {n}: coordinates with no window or screen frame: "
                           f"{script.split(chr(10))[0]!r}")
    return bad


# ---------------------------------------------------------------- the run

INPUT_STEPS = ("drive", "submit", "close")

# What went wrong, collected rather than printed and forgotten: a rehearsal is only
# worth running if it ends in a verdict.
FAILURES = []
REHEARSE = False      # skip the voice, keep every action
BEAT_PAUSE = 0.8      # what a beat costs in a rehearsal
WATCH = True          # stop if the deck leaves fullscreen


class StageLost(Exception):
    """The deck is no longer fullscreen on the recorded monitor."""


def fail(n, msg):
    FAILURES.append(f"slide {n}: {msg}")
    print(f"  [{n}] FAIL: {msg}")


def guard(n, where):
    """Stop the moment the stage is wrong, instead of filming the rest of it.

    The 22:33 take lost its slideshow at 13:04 and carried on for another four
    minutes and five slides. Thirty seconds of checking would have saved all of it."""
    if not WATCH or deck_is_fullscreen():
        return
    fail(n, f"the deck left fullscreen, at {where}")
    raise StageLost(where)


def run_slide(n, steps, beats, dry):
    pending = list(beats)
    for i, step in enumerate(steps):
        nxt = steps[i + 1] if i + 1 < len(steps) else None
        more_input = isinstance(nxt, tuple) and nxt[0] in INPUT_STEPS
        if step == "beat":
            if not pending:
                fail(n, "a beat the deck does not have")
                continue
            text = pending.pop(0)
            if REHEARSE:
                print(f"  [{n}] (beat) {text[:70]}")
                if not dry:
                    time.sleep(BEAT_PAUSE)
            else:
                speak_beat(text, dry, n)
            continue
        kind = step[0]
        if kind == "run":
            print(f"  [{n}] run: {step[1][:90]}")
            if not dry:
                sh(step[1])
        elif kind == "bg":
            print(f"  [{n}] bg:  {step[1][:90]}")
            if not dry:
                sh(step[1], background=True)
        elif kind == "close":
            # This has to come before the generic INPUT_STEPS branch, which matches
            # "close" too and would swallow it. It did: a close step carries
            # (title, script) rather than (script), so the generic branch drove the
            # *title* as a script, close_window never ran, and the viewer and the app
            # window stayed on screen for every slide after. That is what left the
            # report standing in the 22:33 take - not a missed click.
            title, script = step[1], step[2]
            print(f"  [{n}] close: {title}")
            if not dry:
                how = close_window(title, script)
                print(f"  [{n}]   -> {how}")
                if how == "STILL ON SCREEN":
                    fail(n, f"{title} would not close - it would stand in front of "
                            f"every slide after this one")
                if not more_input:
                    park()
        elif kind in INPUT_STEPS:
            print(f"  [{n}] {kind}: {step[1].strip().replace(chr(10), ' | ')[:74]}")
            if not dry:
                if kind == "submit":
                    how = submit(step[1])
                    print(f"  [{n}]   -> {how}")
                    # A card that neither the click nor Return answered is not a
                    # cosmetic miss: the next `clear` dismisses it, so the film shows
                    # a question being waved away while the voice says it was answered,
                    # and the agent blocked on it gets nothing back.
                    if how == "STILL ON SCREEN":
                        fail(n, "a card was not answered - the click and Return both "
                                "left it on screen")
                else:
                    drive(step[1])
                # A pointer that has just clicked is sitting on top of what it
                # clicked, and on this presenter a pointer that has just moved also
                # summons a control strip along the bottom edge. Parking fixes both,
                # but only once the input is finished - not between two fields of the
                # same form, and never on a slide change, which would put that strip
                # on every slide in the video.
                if not more_input and any(k in step[1] for k in ("click", "drag", "move")):
                    park()
        elif kind == "wait_card":
            print(f"  [{n}] wait for the card")
            if not dry and card_geometry() is None:
                fail(n, "no card appeared")
        elif kind == "above":
            print(f"  [{n}] above: {step[1]}")
            if not dry:
                bad = above_stage(step[1])
                if bad:
                    fail(n, f"'{step[1]}' {bad} - the film would show the "
                            f"slideshow where the window should be")
        elif kind == "hold":
            print(f"  [{n}] hold {step[1]}s")
            if not dry:
                time.sleep(min(step[1], 1.0) if REHEARSE else step[1])
        elif kind == "clear":
            print(f"  [{n}] clear")
            if not dry and not clear_cards():
                fail(n, "a card would not close")
        elif kind == "park":
            print(f"  [{n}] park the pointer")
            if not dry:
                park()
        if kind in INPUT_STEPS and not dry:
            guard(n, f"after {kind}: {step[1].splitlines()[0]}")
    if pending:
        fail(n, f"{len(pending)} beat(s) the plan never placed")
        for text in pending:
            if REHEARSE:
                print(f"  [{n}] (beat) {text[:70]}")
            else:
                speak_beat(text, dry, n)


def main():
    global PARK, REHEARSE, WATCH, BEAT_PAUSE
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--dry-run", action="store_true", help="print the timeline, touch nothing")
    ap.add_argument("--rehearse", action="store_true",
                    help="every window, click and card for real, with no narration: "
                         "the run that proves a take will work")
    ap.add_argument("--beat-pause", type=float, default=BEAT_PAUSE,
                    help="seconds a rehearsal spends where a narration beat would be")
    ap.add_argument("--no-watch", action="store_true",
                    help="do not stop when the deck leaves fullscreen (not for a take)")
    ap.add_argument("--from", dest="start", type=int, default=1)
    ap.add_argument("--to", dest="end", type=int, default=len(SLIDES))
    ap.add_argument("--park", default="", help="x,y for the pointer between beats")
    ap.add_argument("--marks", action="store_true",
                    help="set record.sh's in and out marks around the performance")
    args = ap.parse_args()

    if args.park:
        PARK = [int(v) for v in args.park.split(",")]
    REHEARSE = args.rehearse
    BEAT_PAUSE = args.beat_pause
    WATCH = not args.no_watch

    beats = beats_from_deck()
    total = sum(len(b) for b in beats.values())
    words = sum(len(" ".join(b).split()) for b in beats.values())
    print(f"{len(beats)} slides, {total} beats, {words} words "
          f"(about {words / 2.9 / 60:.1f} min of narration)"
          f"{'  [REHEARSAL: no voice]' if REHEARSE else ''}\n")

    # The plan itself, before anything moves. A scroll with no move in front of it is
    # the mistake that cost the 22:33 take.
    problems = lint_plan()
    for p in problems:
        print(f"LINT: {p}")
    if problems and not args.dry_run:
        sys.exit("the plan is unsafe to perform; fix the steps above")

    # A beat the deck does not have is silence where a sentence was meant to be, and
    # one the plan does not ask for arrives after the visuals it belongs to. Both are
    # only visible in the finished video, so they are reported before the camera rolls.
    if len(beats) != len(SLIDES):
        print(f"WARNING: the deck has {len(beats)} slides, the plan has {len(SLIDES)}")
    for n, steps in sorted(SLIDES.items()):
        want = sum(1 for s in steps if s == "beat")
        have = len(beats.get(n, []))
        if want != have:
            print(f"WARNING: slide {n} narrates {have} beat(s), the plan places {want}")

    if not args.dry_run:
        if WATCH and not deck_is_fullscreen():
            sys.exit("the deck is not fullscreen on the recorded monitor: start its "
                     "slideshow first, or pass --no-watch")

    if not args.dry_run and args.start == 1:
        # Park first, before anything opens. Every window AgentBox places goes to the
        # monitor the pointer is on, so the pointer decides which screen the demo
        # happens on - and a take records one monitor.
        park()
        drive("key home\n")  # the deck starts where the narration does
        time.sleep(0.8)

    # The marks belong in here rather than in the hands of whoever starts this. The
    # first rehearsal put them in two separate shell commands and the thirteen
    # seconds between those commands landed inside the video, as a title slide in
    # silence. Set from here, the cut begins on the first word and ends on the last.
    def mark(which):
        if args.marks and not args.dry_run:
            subprocess.run([RECORD, "mark", which], capture_output=True, text=True)

    began = time.time()
    timing = []
    mark("in")
    try:
        for n in range(args.start, args.end + 1):
            at = time.time() - began
            print(f"--- slide {n}  (+{at:.0f}s)")
            if not args.dry_run:
                guard(n, "the start of the slide")
            run_slide(n, SLIDES.get(n, ["beat"]), beats.get(n, []), args.dry_run)
            timing.append((n, at, time.time() - began - at))
            if n < args.end:
                if not args.dry_run:
                    drive("key right\n")
                    time.sleep(0.4)
                print(f"  [{n}] advance")
    except StageLost as e:
        print(f"\nSTOPPED: {e}. Nothing after this point would have been usable.")
    mark("out")
    print(f"\ndone in {time.time() - began:.0f}s")

    if REHEARSE:
        print("\n--- rehearsal report")
        print(f"{'slide':>5}  {'starts':>7}  {'takes':>6}")
        for n, at, dur in timing:
            print(f"{n:>5}  {at:>6.0f}s  {dur:>5.1f}s")
        action = sum(d for _, _, d in timing)
        spoken = words / 2.9
        paused = total * BEAT_PAUSE
        print(f"\n{action:.0f}s of actions and holds, of which {paused:.0f}s stood in "
              f"for narration.")
        print(f"With the real narration ({spoken:.0f}s) that is a take of about "
              f"{(action - paused + spoken) / 60:.0f} minutes.")
    if FAILURES:
        print(f"\n{len(FAILURES)} problem(s) - do not record until these are gone:")
        for f in FAILURES:
            print(f"  {f}")
        sys.exit(1)
    print("\nno problems found" if REHEARSE else "")


if __name__ == "__main__":
    main()
